import streamlit as st
from pptx import Presentation
import io
import base64

# アプリのタイトルと作成者情報
st.title("PowerPoint Slide Viewer")
st.caption("Created by Dit-Lab.(Daiki Ito)")

# アプリの説明
st.markdown("""
## **概要**
このウェブアプリケーションでは、PowerPointファイルをアップロードし、スライドとメモを表示することができます。
""")

# ファイルアップローダー
uploaded_file = st.file_uploader("PowerPointファイルをアップロードしてください", type=['pptx'])

if uploaded_file is not None:
    # PowerPointファイルを読み込む
    prs = Presentation(uploaded_file)
    
    # スライドの選択
    total_slides = len(prs.slides)
    slide_number = st.slider("スライドを選択", 1, total_slides, 1)
    
    # 選択したスライドとそのメモを表示
    current_slide = prs.slides[slide_number - 1]
    
    # スライドを画像として保存
    image_stream = io.BytesIO()
    current_slide.save(image_stream, "PNG")
    image_stream.seek(0)
    image_data = base64.b64encode(image_stream.getvalue()).decode()
    
    # スライドを画像として表示
    st.image(f"data:image/png;base64,{image_data}", use_column_width=True)
    
    # スライドのメモを表示
    st.subheader("スライドメモ")
    notes_slide = current_slide.notes_slide
    if notes_slide and notes_slide.notes_text_frame.text:
        st.write(notes_slide.notes_text_frame.text)
    else:
        st.write("このスライドにメモはありません。")

# コピーライト情報
st.markdown("---")
st.subheader('© 2022-2024 Dit-Lab.(Daiki Ito). All Rights Reserved.')
st.write("PowerPoint Slide Viewer: Making presentations accessible")
st.write("Democratizing presentation sharing, everywhere.")